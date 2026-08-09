from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1F, 0x4E, 0x79)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)

FOOTER_TEXT = "consultor-ia-pppm  ·  Geovane Virmecati  ·  Eixo Estratégico"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def set_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_accent_bar(slide, color=ORANGE, top=Inches(0.5), left=Inches(0.6), width=Inches(1.2), height=Inches(0.08)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_footer(slide, page_num=None, total=10):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(9), Inches(0.35))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = FOOTER_TEXT
    r.font.name = "Helvetica"
    r.font.size = Pt(9)
    r.font.color.rgb = GREY
    if page_num is not None:
        tb2 = slide.shapes.add_textbox(Inches(11.7), Inches(7.05), Inches(1.2), Inches(0.35))
        tf2 = tb2.text_frame
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = f"{page_num} / {total}"
        r2.font.name = "Helvetica"
        r2.font.size = Pt(9)
        r2.font.color.rgb = GREY


def add_title(slide, text, top=Inches(0.75), left=Inches(0.6), width=Inches(12), size=32, color=NAVY, bold=True):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.9))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Helvetica"
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def add_subtitle(slide, text, top=Inches(1.55), left=Inches(0.6), width=Inches(12), size=16, color=GREY):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Helvetica"
    r.font.size = Pt(size)
    r.font.color.rgb = color
    return tb


def add_bullets(slide, bullets, top=Inches(2.4), left=Inches(0.9), width=Inches(11.5), size=18, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = f"•  {b}"
        r.font.name = "Helvetica"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_quote_box(slide, quote, source=None, top=Inches(4.7), left=Inches(1.5), width=Inches(10.3), height=Inches(1.6)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = quote
    r.font.name = "Helvetica"
    r.font.size = Pt(20)
    r.font.color.rgb = WHITE
    r.font.bold = True
    if source:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        r2 = p2.add_run()
        r2.text = source
        r2.font.name = "Helvetica"
        r2.font.size = Pt(11)
        r2.font.color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
        r2.font.italic = True


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================================
# SLIDE 1 — Capa
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)
# barra laranja lateral
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
side.fill.solid(); side.fill.fore_color.rgb = ORANGE; side.line.fill.background()

# título
tb = s.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Do método à ferramenta"
r.font.name = "Helvetica"; r.font.size = Pt(52); r.font.color.rgb = WHITE; r.font.bold = True

tb = s.shapes.add_textbox(Inches(1), Inches(3.4), Inches(11), Inches(1))
tf = tb.text_frame
p = tf.paragraphs[0]; r = p.add_run()
r.text = "consultor-ia-pppm"
r.font.name = "Helvetica"; r.font.size = Pt(38); r.font.color.rgb = ORANGE; r.font.bold = True

tb = s.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11), Inches(1))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Um app que materializa o método da Aula 1 do Prof. Dr. José Bezerra"
r.font.name = "Helvetica"; r.font.size = Pt(18); r.font.color.rgb = LIGHT

# autor
tb = s.shapes.add_textbox(Inches(1), Inches(6.4), Inches(11), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Geovane Virmecati  ·  Eixo Estratégico  ·  Formação de Consultores em IA aplicada ao PPPM — BSBr"
r.font.name = "Helvetica"; r.font.size = Pt(12); r.font.color.rgb = LIGHT

add_notes(s, """Abertura curta. Não me alonga.

'Bom dia a todos. Meu nome é Geovane Virmecati. Nos próximos 10 minutos vocês vão ver o método que o professor acabou de apresentar transformado em ferramenta. E o mais importante: vão ver COMO essa ferramenta foi construída — porque o processo é exatamente o que o curso está formando.'""")


# ============================================================
# SLIDE 2 — A provocação
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_accent_bar(s)
add_title(s, "A provocação")
add_subtitle(s, "Vocês acabaram de ver um método. Agora vão ver ele virado ferramenta em 5 minutos.")

add_quote_box(s, '"IA sem método vira ferramenta. IA com método vira valor."', "Prof. Dr. José Bezerra · Aula 1 · slide 22", top=Inches(3.2), height=Inches(2))

add_footer(s, 2)
add_notes(s, """Deixa a frase pousar. Olhe para o professor quando disser isso.

'Essa frase do professor é o eixo da apresentação. O que vocês vão ver a seguir é a tentativa mais literal possível de honrar esse princípio: usar IA como método, não como ferramenta.'""")


# ============================================================
# SLIDE 3 — O que o app faz
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_accent_bar(s)
add_title(s, "O que o app faz")
add_subtitle(s, "Fluxo em 5 etapas — o mesmo método da Empresa Alfa, com os dados do participante.")

# 5 caixas em linha
steps = [
    ("1", "Contexto", "empresa, porte,\nprojetos, PMO"),
    ("2", "Diagnóstico", "5 dimensões\n× 0-6 pontos"),
    ("3", "Mapa 5 Blocos", "Contexto, Dor,\nDados, Riscos, Valor"),
    ("4", "3 Pilotos", "recomendação\ndeterminística"),
    ("5", "PDF", "entregável\nexecutivo"),
]
box_w = Inches(2.3); box_h = Inches(2.6); gap = Inches(0.15)
total_w = 5 * box_w + 4 * gap
start_left = (SW - total_w) / 2

for i, (num, title, desc) in enumerate(steps):
    left = start_left + i * (box_w + gap)
    top = Inches(2.6)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = ORANGE; r.font.name = "Helvetica"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = title; r2.font.size = Pt(15); r2.font.bold = True; r2.font.color.rgb = NAVY; r2.font.name = "Helvetica"
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(10)
    r3 = p3.add_run(); r3.text = desc; r3.font.size = Pt(11); r3.font.color.rgb = DARK; r3.font.name = "Helvetica"

# tempo total
tb = s.shapes.add_textbox(Inches(0.6), Inches(6.1), Inches(12), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Tempo total do participante: ~10 minutos"
r.font.name = "Helvetica"; r.font.size = Pt(14); r.font.color.rgb = ORANGE; r.font.italic = True

add_footer(s, 3)
add_notes(s, """Aponta cada caixa. Não leia em voz alta — os alunos leem sozinhos.

'É literalmente o Mapa 5 Blocos da Empresa Alfa que o professor mostrou. A diferença é que agora cada aluno preenche com o CONTEXTO DELE, e sai da sala com o PDF já pronto para levar para a diretoria.'""")


# ============================================================
# SLIDE 4 — Demo ao vivo
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)

# título grande centrado
tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Demo ao vivo"
r.font.name = "Helvetica"; r.font.size = Pt(60); r.font.color.rgb = WHITE; r.font.bold = True

tb = s.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Vou rodar aqui — na minha máquina, offline, sem chamar API."
r.font.name = "Helvetica"; r.font.size = Pt(20); r.font.color.rgb = LIGHT; r.font.italic = True

tb = s.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "http://localhost:8512"
r.font.name = "Courier New"; r.font.size = Pt(24); r.font.color.rgb = ORANGE; r.font.bold = True

add_notes(s, """Este slide é só para o público entender que agora começa a demo. Não fica muito tempo aqui — abre o app no navegador e roda o caso ao vivo.

Sequência da demo:
1. Preenche contexto (30s) — usa um nome fictício ou pede um voluntário da sala
2. Diagnóstico com 5 sliders (1min) — pontuações que caem no Nível 3
3. Mapa 5 Blocos (2min) — cola textos preparados OU escreve ao vivo baseado no caso da sala
4. Aba 4 mostra os 3 pilotos aparecendo automaticamente
5. Gera o PDF na aba 5 e projeta

Total demo: ~5 minutos. Se der bug, pule para o slide 5 (não pare).""")


# ============================================================
# SLIDE 5 — Pipeline formal
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_accent_bar(s)
add_title(s, "Como foi construído — pipeline formal")
add_subtitle(s, 'Não foi "prompt aleatório para IA". Foi processo com etapas, gates e validação humana.')

# 8 etapas em 2 linhas de 4
etapas = [
    ("1", "Backlog"),
    ("2", "Triagem"),
    ("3", "Especificação"),
    ("4", "Arquitetura"),
    ("5", "Construção"),
    ("6", "Validação"),
    ("7", "Deploy"),
    ("8", "Documentação"),
]
box_w = Inches(2.75); box_h = Inches(1.5); gap_h = Inches(0.15); gap_v = Inches(0.35)
row_total = 4 * box_w + 3 * gap_h
row_left = (SW - row_total) / 2

for i, (num, nome) in enumerate(etapas):
    row = i // 4; col = i % 4
    left = row_left + col * (box_w + gap_h)
    top = Inches(2.5) + row * (box_h + gap_v)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY; box.line.width = Pt(1)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"{num}. "; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = ORANGE; r.font.name = "Helvetica"
    r2 = p.add_run(); r2.text = nome; r2.font.size = Pt(18); r2.font.bold = True; r2.font.color.rgb = NAVY; r2.font.name = "Helvetica"

# nota inferior
tb = s.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Cada etapa produziu artefato — arquitetura em 17 KB de markdown ANTES da primeira linha de código."
r.font.name = "Helvetica"; r.font.size = Pt(13); r.font.color.rgb = GREY; r.font.italic = True

add_footer(s, 5)
add_notes(s, """Devagar aqui. Este é o coração da apresentação.

'Reparem: das 8 etapas, só a 5 é CODAR. As outras 7 são MÉTODO. O professor falou hoje que consultor é quem pensa antes, não quem executa mais rápido. Foi exatamente isso que eu tentei fazer.'

Ponto forte para a diretoria: 17KB de arquitetura em markdown ANTES do primeiro `def`.""")


# ============================================================
# SLIDE 6 — Agentes especializados
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_accent_bar(s)
add_title(s, "Agentes especializados, não IA generalista")
add_subtitle(s, "Cada agente com uma responsabilidade única + validação humana (HITL) entre etapas.")

agentes = [
    ("architect", "Desenhou a arquitetura", "estrutura, módulos, schemas, regras"),
    ("builder", "Implementou o código", "24 arquivos Python + 3 JSON"),
    ("vault-writer", "Documentou no Obsidian", "nota estruturada com frontmatter"),
    ("librarian", "Validou a documentação", "checou índice, tags, links"),
]

for i, (nome, papel, detalhe) in enumerate(agentes):
    top = Inches(2.5) + i * Inches(1.0)
    # bloco esquerdo (nome do agente)
    box1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), top, Inches(2.8), Inches(0.85))
    box1.fill.solid(); box1.fill.fore_color.rgb = NAVY; box1.line.fill.background()
    tf1 = box1.text_frame
    tf1.margin_left = Inches(0.2); tf1.margin_right = Inches(0.2)
    tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf1.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = nome
    r1.font.name = "Courier New"; r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = WHITE

    # papel
    tb = s.shapes.add_textbox(Inches(4.0), top, Inches(4.5), Inches(0.85))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = papel
    r.font.name = "Helvetica"; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = DARK

    # detalhe
    tb = s.shapes.add_textbox(Inches(8.6), top, Inches(4.5), Inches(0.85))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = detalhe
    r.font.name = "Helvetica"; r.font.size = Pt(14); r.font.color.rgb = GREY; r.font.italic = True

# frase âncora
tb = s.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Nenhum agente sabia o que o outro estava fazendo. Escopo restrito = auditabilidade."
r.font.name = "Helvetica"; r.font.size = Pt(13); r.font.color.rgb = ORANGE; r.font.italic = True

add_footer(s, 6)
add_notes(s, """'Não pedi para uma IA fazer tudo. Distribuí em 4 agentes, cada um com UMA função. Isso é literalmente separação de responsabilidades — princípio básico de engenharia de software aplicado à orquestração de IA.'

Se alguém perguntar sobre HITL, responde: 'Entre cada agente eu revisei manualmente o output antes do próximo entrar. IA sugere, humano valida, próximo agente executa. Isso é HITL de verdade — não é botão de aprovar no fim.'""")


# ============================================================
# SLIDE 7 — Zero LLM em produção
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_accent_bar(s, color=ORANGE)
add_title(s, "A decisão contraintuitiva")

# número gigante
tb = s.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(2))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "ZERO"
r.font.name = "Helvetica"; r.font.size = Pt(120); r.font.color.rgb = ORANGE; r.font.bold = True

tb = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "chamadas de LLM em produção"
r.font.name = "Helvetica"; r.font.size = Pt(24); r.font.color.rgb = NAVY; r.font.bold = True

# 3 justificativas
justif = ["AUDITÁVEL", "GRATUITO", "RASTREÁVEL"]
w = Inches(3.2); gap = Inches(0.4); total = 3 * w + 2 * gap
start = (SW - total) / 2
for i, txt in enumerate(justif):
    left = start + i * (w + gap)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.3), w, Inches(0.8))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.name = "Helvetica"; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Recomendador de pilotos é regra matemática, não modelo estatístico. Diretoria consegue auditar."
r.font.name = "Helvetica"; r.font.size = Pt(13); r.font.color.rgb = GREY; r.font.italic = True

add_footer(s, 7)
add_notes(s, """Este slide provoca. Pause aqui.

'Alguém deve estar pensando: "Mas se é um app de IA, por que não usa IA para decidir os pilotos?" Boa pergunta.

RESPOSTA: porque a aula ensina que IA sem governança vira caixa-preta. Se cada execução chamasse LLM, o app seria menos consultivo, não mais. Diretoria pediria: "prove que a recomendação está certa". Com regra determinística, eu abro o código e mostro. Com LLM, eu digo "confia".

IA foi usada NA CONSTRUÇÃO. NÃO NA OPERAÇÃO. Essa é a diferença entre consultor e vendedor de ferramenta.'""")


# ============================================================
# SLIDE 8 — Camada meta (3 frases da aula)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_accent_bar(s)
add_title(s, "3 frases da aula → 3 decisões do projeto")
add_subtitle(s, "A camada meta: o próprio app é caso vivo do que a aula ensina.")

pares = [
    ("IA com método vira valor.", "Pipeline formal + 28 testes automatizados", "→"),
    ("Governança e HITL são o gargalo.", "Agentes com escopo restrito + validação humana entre etapas", "→"),
    ("Diagnóstico revela oportunidade.", "App respeita o espaço que o professor deixou para próximas aulas", "→"),
]

for i, (frase, decisao, seta) in enumerate(pares):
    top = Inches(2.4) + i * Inches(1.3)
    # frase à esquerda (aula)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), top, Inches(5.4), Inches(1.0))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f'"{frase}"'
    r.font.name = "Helvetica"; r.font.size = Pt(14); r.font.color.rgb = WHITE; r.font.italic = True

    # seta
    tb = s.shapes.add_textbox(Inches(6.35), top, Inches(0.5), Inches(1.0))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "→"
    r.font.name = "Helvetica"; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = ORANGE

    # decisão à direita (app)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), top, Inches(5.6), Inches(1.0))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = ORANGE; box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = decisao
    r.font.name = "Helvetica"; r.font.size = Pt(13); r.font.color.rgb = DARK; r.font.bold = True

add_footer(s, 8)
add_notes(s, """Este é O slide da apresentação. Fala devagar.

'Cada frase que o professor disse hoje virou uma decisão de arquitetura deste app. Não foi coincidência. Foi tentativa DELIBERADA de honrar o método aplicando ele em cima de si mesmo.'

Última linha, olhando para o professor: 'Se o app funciona, é porque o método funciona.'""")


# ============================================================
# SLIDE 9 — Números
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_accent_bar(s)
add_title(s, "Números do projeto")
add_subtitle(s, "Reproduzível. Barato. Rastreável.")

numeros = [
    ("~1h", "do 'começa' ao\nPDF funcionando"),
    ("12", "arquivos Python\n+ 3 JSON"),
    ("28/28", "testes\nautomatizados"),
    ("0", "chamadas de IA\npor execução"),
    ("R$ 0", "custo\npor uso"),
    ("100%", "código aberto\ndeterminístico"),
]

box_w = Inches(3.9); box_h = Inches(1.8); gap_h = Inches(0.15); gap_v = Inches(0.2)
row_total = 3 * box_w + 2 * gap_h
row_left = (SW - row_total) / 2

for i, (n, label) in enumerate(numeros):
    row = i // 3; col = i % 3
    left = row_left + col * (box_w + gap_h)
    top = Inches(2.5) + row * (box_h + gap_v)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = NAVY; box.line.width = Pt(1)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n
    r.font.name = "Helvetica"; r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = ORANGE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = label
    r2.font.name = "Helvetica"; r2.font.size = Pt(12); r2.font.color.rgb = DARK

# rodapé
tb = s.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Repositório: github.com/geovime1977/consultor-ia-pppm  ·  Backup: OneDrive Eixo Estratégico"
r.font.name = "Helvetica"; r.font.size = Pt(11); r.font.color.rgb = GREY

add_footer(s, 9)
add_notes(s, """Vá lendo os números um por um. Deixe o público absorver.

Ponto forte para provocação:
'R$ 0 por uso. Zero. Cada aluno na sala pode instalar isso no notebook dele e usar quantas vezes quiser sem pagar centavo. Diretoria adora ROI infinito.'""")


# ============================================================
# SLIDE 10 — Fechamento
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
side.fill.solid(); side.fill.fore_color.rgb = ORANGE; side.line.fill.background()

# frase-âncora
tb = s.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11), Inches(2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "IA foi usada NA construção."
r.font.name = "Helvetica"; r.font.size = Pt(38); r.font.color.rgb = WHITE; r.font.bold = True
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
r2 = p2.add_run(); r2.text = "Não na operação."
r2.font.name = "Helvetica"; r2.font.size = Pt(38); r2.font.color.rgb = ORANGE; r2.font.bold = True

tb = s.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "É essa a diferença que este curso está formando."
r.font.name = "Helvetica"; r.font.size = Pt(18); r.font.color.rgb = LIGHT; r.font.italic = True

# créditos
tb = s.shapes.add_textbox(Inches(1), Inches(6.0), Inches(11), Inches(0.9))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Método: Prof. Dr. José Bezerra · BSBr  |  Base: PMI Standard for AI in PPPM (2026)"
r.font.name = "Helvetica"; r.font.size = Pt(11); r.font.color.rgb = LIGHT
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "github.com/geovime1977/consultor-ia-pppm  ·  Geovane Virmecati · Eixo Estratégico"
r2.font.name = "Helvetica"; r2.font.size = Pt(11); r2.font.color.rgb = LIGHT

add_notes(s, """Fechamento curto e firme.

'Essa é a mensagem que eu queria deixar. IA sem método vira ferramenta. IA com método vira valor. Este app é a tentativa mais literal possível de honrar o método do professor Bezerra — usando ele em cima de si mesmo.

Obrigado pela atenção. Se algum de vocês quiser instalar na própria máquina, o professor tem o pacote com um instalador de duplo-clique. Zero custo, zero dependência de servidor.

Perguntas?'""")


# ============================================================
# SALVAR
# ============================================================
out = "/Users/virmecati/projetos/consultor-ia-pppm/docs/DECK-PALESTRA.pptx"
prs.save(out)
import os
print(f"PPTX gerado: {out} ({os.path.getsize(out)/1024:.1f} KB, {len(prs.slides)} slides)")
